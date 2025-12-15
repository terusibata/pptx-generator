"""
PowerPoint生成サービス
テンプレートとAI出力のJSONからPPTXを生成
"""
from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR, MSO_LINE_DASH_STYLE
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches

from app.schemas import (
    BulletContent,
    ChartContent,
    ChartType,
    ConnectorContent,
    ContentType,
    ImageContent,
    ParagraphContent,
    PlaceholderType,
    ShapeContent,
    ShapeStyle,
    ShapeType,
    SlideContent,
    SlideDefinition,
    TableContent,
    TemplateMeta,
    TextBoxContent,
    TextContent,
    TextStyle,
)

if TYPE_CHECKING:
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.placeholder import SlidePlaceholder
    from pptx.slide import Slide, SlideLayout
    from pptx.text.text import TextFrame


# =============================================================================
# 定数・マッピング
# =============================================================================

ALIGNMENT_MAP: dict[str, int] = {
    "LEFT": PP_ALIGN.LEFT,
    "CENTER": PP_ALIGN.CENTER,
    "RIGHT": PP_ALIGN.RIGHT,
    "JUSTIFY": PP_ALIGN.JUSTIFY,
}

CHART_TYPE_MAP: dict[ChartType, int] = {
    ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
    ChartType.COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
    ChartType.LINE: XL_CHART_TYPE.LINE,
    ChartType.PIE: XL_CHART_TYPE.PIE,
    ChartType.DOUGHNUT: XL_CHART_TYPE.DOUGHNUT,
    ChartType.AREA: XL_CHART_TYPE.AREA,
}

THEME_COLOR_MAP: dict[str, int] = {
    "DARK_1": MSO_THEME_COLOR.DARK_1,
    "LIGHT_1": MSO_THEME_COLOR.LIGHT_1,
    "DARK_2": MSO_THEME_COLOR.DARK_2,
    "LIGHT_2": MSO_THEME_COLOR.LIGHT_2,
    "ACCENT_1": MSO_THEME_COLOR.ACCENT_1,
    "ACCENT_2": MSO_THEME_COLOR.ACCENT_2,
    "ACCENT_3": MSO_THEME_COLOR.ACCENT_3,
    "ACCENT_4": MSO_THEME_COLOR.ACCENT_4,
    "ACCENT_5": MSO_THEME_COLOR.ACCENT_5,
    "ACCENT_6": MSO_THEME_COLOR.ACCENT_6,
}

SHAPE_TYPE_MAP: dict[ShapeType, int] = {
    ShapeType.RECTANGLE: MSO_SHAPE.RECTANGLE,
    ShapeType.ROUNDED_RECTANGLE: MSO_SHAPE.ROUNDED_RECTANGLE,
    ShapeType.OVAL: MSO_SHAPE.OVAL,
    ShapeType.TRIANGLE: MSO_SHAPE.ISOSCELES_TRIANGLE,
    ShapeType.RIGHT_TRIANGLE: MSO_SHAPE.RIGHT_TRIANGLE,
    ShapeType.DIAMOND: MSO_SHAPE.DIAMOND,
    ShapeType.PENTAGON: MSO_SHAPE.PENTAGON,
    ShapeType.HEXAGON: MSO_SHAPE.HEXAGON,
    ShapeType.ARROW_RIGHT: MSO_SHAPE.RIGHT_ARROW,
    ShapeType.ARROW_LEFT: MSO_SHAPE.LEFT_ARROW,
    ShapeType.ARROW_UP: MSO_SHAPE.UP_ARROW,
    ShapeType.ARROW_DOWN: MSO_SHAPE.DOWN_ARROW,
    ShapeType.CHEVRON: MSO_SHAPE.CHEVRON,
    ShapeType.STAR_5_POINT: MSO_SHAPE.STAR_5_POINT,
    ShapeType.STAR_6_POINT: MSO_SHAPE.STAR_6_POINT,
    ShapeType.CALLOUT_RECTANGULAR: MSO_SHAPE.RECTANGULAR_CALLOUT,
    ShapeType.CALLOUT_ROUNDED_RECTANGULAR: MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
    ShapeType.CALLOUT_OVAL: MSO_SHAPE.OVAL_CALLOUT,
    ShapeType.CALLOUT_CLOUD: MSO_SHAPE.CLOUD_CALLOUT,
    ShapeType.CURVED_RIGHT_ARROW: MSO_SHAPE.CURVED_RIGHT_ARROW,
    ShapeType.CURVED_LEFT_ARROW: MSO_SHAPE.CURVED_LEFT_ARROW,
    ShapeType.CURVED_UP_ARROW: MSO_SHAPE.CURVED_UP_ARROW,
    ShapeType.CURVED_DOWN_ARROW: MSO_SHAPE.CURVED_DOWN_ARROW,
    ShapeType.BLOCK_ARC: MSO_SHAPE.BLOCK_ARC,
    ShapeType.DONUT: MSO_SHAPE.DONUT,
    ShapeType.HEART: MSO_SHAPE.HEART,
    ShapeType.LIGHTNING_BOLT: MSO_SHAPE.LIGHTNING_BOLT,
    ShapeType.SUN: MSO_SHAPE.SUN,
    ShapeType.MOON: MSO_SHAPE.MOON,
    ShapeType.CLOUD: MSO_SHAPE.CLOUD,
    ShapeType.FLOWCHART_PROCESS: MSO_SHAPE.FLOWCHART_PROCESS,
    ShapeType.FLOWCHART_DECISION: MSO_SHAPE.FLOWCHART_DECISION,
    ShapeType.FLOWCHART_TERMINATOR: MSO_SHAPE.FLOWCHART_TERMINATOR,
    ShapeType.FLOWCHART_DATA: MSO_SHAPE.FLOWCHART_DATA,
    ShapeType.FLOWCHART_CONNECTOR: MSO_SHAPE.FLOWCHART_CONNECTOR,
}

LINE_DASH_MAP: dict[str, int] = {
    "solid": MSO_LINE_DASH_STYLE.SOLID,
    "dash": MSO_LINE_DASH_STYLE.DASH,
    "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
}


# =============================================================================
# ヘルパー関数
# =============================================================================

def _parse_rgb_color(color_str: str) -> RGBColor:
    """RGB文字列をRGBColorオブジェクトに変換"""
    color_str = color_str.lstrip("#")
    if len(color_str) != 6:
        raise ValueError(f"Invalid color format: {color_str}")
    r = int(color_str[0:2], 16)
    g = int(color_str[2:4], 16)
    b = int(color_str[4:6], 16)
    return RGBColor(r, g, b)


def _apply_text_style(run, style: TextStyle | None) -> None:
    """テキストランにスタイルを適用"""
    if style is None:
        return
    
    if style.bold:
        run.font.bold = True
    if style.italic:
        run.font.italic = True
    if style.underline:
        run.font.underline = True
    if style.font_size:
        run.font.size = Pt(style.font_size)
    if style.font_name:
        run.font.name = style.font_name
    if style.color:
        run.font.color.rgb = _parse_rgb_color(style.color)
    elif style.theme_color and style.theme_color in THEME_COLOR_MAP:
        run.font.color.theme_color = THEME_COLOR_MAP[style.theme_color]


def _find_placeholder(
    slide: "Slide",
    name: str | None = None,
    idx: int | None = None,
    ph_type: PlaceholderType | None = None,
) -> "SlidePlaceholder | None":
    """
    スライドからプレースホルダーを検索
    優先順位: idx > name > ph_type
    """
    for shape in slide.placeholders:
        if not hasattr(shape, "placeholder_format"):
            continue
        
        ph_format = shape.placeholder_format
        
        # インデックスで検索
        if idx is not None and ph_format.idx == idx:
            return shape
        
        # 名前で検索（部分一致）
        if name is not None:
            shape_name = shape.name or ""
            if name.lower() in shape_name.lower() or shape_name.lower() in name.lower():
                return shape
        
        # タイプで検索
        if ph_type is not None:
            type_map = {
                PlaceholderType.TITLE: PP_PLACEHOLDER.TITLE,
                PlaceholderType.CENTER_TITLE: PP_PLACEHOLDER.CENTER_TITLE,
                PlaceholderType.SUBTITLE: PP_PLACEHOLDER.SUBTITLE,
                PlaceholderType.BODY: PP_PLACEHOLDER.BODY,
            }
            if ph_type in type_map and ph_format.type == type_map[ph_type]:
                return shape
    
    return None


def _get_all_layouts(prs: Presentation) -> list["SlideLayout"]:
    """
    全スライドマスターから全レイアウトを取得

    prs.slide_layouts は最初のスライドマスターのレイアウトのみを返すため、
    複数のスライドマスターがある場合はこの関数を使用する
    """
    layouts = []
    for master in prs.slide_masters:
        layouts.extend(master.slide_layouts)
    return layouts


def _find_layout_by_name_or_index(
    prs: Presentation,
    name: str | None = None,
    index: int | None = None,
    meta: TemplateMeta | None = None,
) -> "SlideLayout":
    """
    レイアウトを名前またはインデックスで検索
    """
    # 全スライドマスターから全レイアウトを取得
    all_layouts = _get_all_layouts(prs)

    # インデックス指定の場合
    if index is not None:
        if 0 <= index < len(all_layouts):
            return all_layouts[index]
        raise ValueError(f"Layout index {index} out of range (0-{len(all_layouts)-1})")

    # 名前指定の場合
    if name is not None:
        # エイリアス検索
        if meta and meta.custom_config and meta.custom_config.layout_aliases:
            if name in meta.custom_config.layout_aliases:
                alias_index = meta.custom_config.layout_aliases[name]
                if 0 <= alias_index < len(all_layouts):
                    return all_layouts[alias_index]

        # 名前で検索（完全一致 → 部分一致）
        for layout in all_layouts:
            if layout.name == name:
                return layout

        # 部分一致
        name_lower = name.lower()
        for layout in all_layouts:
            if name_lower in (layout.name or "").lower():
                return layout

        raise ValueError(f"Layout not found: {name}")

    # デフォルトは最初のレイアウト
    return all_layouts[0] if all_layouts else prs.slide_layouts[0]


# =============================================================================
# コンテンツ挿入関数
# =============================================================================

def _insert_text_content(
    text_frame: "TextFrame",
    content: TextContent,
    clear_existing: bool = True,
) -> None:
    """テキストコンテンツをテキストフレームに挿入"""
    if clear_existing:
        text_frame.clear()
    
    for i, para_content in enumerate(content.paragraphs):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        # テキスト設定
        run = p.add_run()
        run.text = para_content.text
        
        # スタイル適用
        _apply_text_style(run, para_content.style)
        
        # 段落プロパティ
        if para_content.bullet:
            p.level = para_content.level
        
        if para_content.alignment and para_content.alignment in ALIGNMENT_MAP:
            p.alignment = ALIGNMENT_MAP[para_content.alignment]


def _insert_bullet_content(
    text_frame: "TextFrame",
    content: BulletContent,
    clear_existing: bool = True,
) -> None:
    """箇条書きコンテンツをテキストフレームに挿入"""
    if clear_existing:
        text_frame.clear()
    
    for i, item in enumerate(content.items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.level = 0  # 箇条書きレベル
        
        if isinstance(item, str):
            run = p.add_run()
            run.text = item
            _apply_text_style(run, content.style)
        else:
            # ParagraphContent の場合
            run = p.add_run()
            run.text = item.text
            _apply_text_style(run, item.style or content.style)
            p.level = item.level


def _insert_simple_text(
    text_frame: "TextFrame",
    text: str,
    clear_existing: bool = True,
) -> None:
    """単純なテキストを挿入"""
    if clear_existing:
        text_frame.clear()
    
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text


def _insert_simple_bullets(
    text_frame: "TextFrame",
    items: list[str],
    clear_existing: bool = True,
) -> None:
    """単純な箇条書きリストを挿入"""
    if clear_existing:
        text_frame.clear()
    
    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.level = 0
        run = p.add_run()
        run.text = item


def _insert_table(
    slide: "Slide",
    placeholder: "SlidePlaceholder",
    content: TableContent,
) -> None:
    """テーブルコンテンツを挿入"""
    rows_count = len(content.rows)
    if content.headers:
        rows_count += 1
    
    cols_count = len(content.rows[0]) if content.rows else len(content.headers or [])
    
    # テーブルを作成
    table = slide.shapes.add_table(
        rows_count,
        cols_count,
        placeholder.left,
        placeholder.top,
        placeholder.width,
        placeholder.height,
    ).table
    
    row_offset = 0
    
    # ヘッダー行
    if content.headers:
        for col_idx, header in enumerate(content.headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            if content.style and content.style.header_bg_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _parse_rgb_color(content.style.header_bg_color)
        row_offset = 1
    
    # データ行
    for row_idx, row in enumerate(content.rows):
        for col_idx, cell_data in enumerate(row):
            cell = table.cell(row_idx + row_offset, col_idx)
            if isinstance(cell_data, str):
                cell.text = cell_data
            else:
                cell.text = cell_data.text


def _insert_chart(
    slide: "Slide",
    placeholder: "SlidePlaceholder",
    content: ChartContent,
) -> None:
    """チャートコンテンツを挿入"""
    chart_data = CategoryChartData()
    chart_data.categories = content.categories
    
    for series in content.series:
        chart_data.add_series(series.name, series.values)
    
    chart_type = CHART_TYPE_MAP.get(content.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    chart = slide.shapes.add_chart(
        chart_type,
        placeholder.left,
        placeholder.top,
        placeholder.width,
        placeholder.height,
        chart_data,
    ).chart
    
    if content.title:
        chart.has_title = True
        chart.chart_title.text_frame.text = content.title
    
    chart.has_legend = content.show_legend


def _insert_image(
    slide: "Slide",
    placeholder: "SlidePlaceholder",
    content: ImageContent,
) -> None:
    """画像コンテンツを挿入"""
    image_path = Path(content.source)

    if not image_path.exists():
        # URL の場合は後で対応（ここではスキップ）
        return

    slide.shapes.add_picture(
        str(image_path),
        placeholder.left,
        placeholder.top,
        placeholder.width,
        placeholder.height,
    )


def _apply_shape_style(shape, style: ShapeStyle | None) -> None:
    """図形にスタイルを適用"""
    if style is None:
        return

    # 塗りつぶし色
    if style.fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _parse_rgb_color(style.fill_color)

    # 線の色
    if style.line_color:
        shape.line.color.rgb = _parse_rgb_color(style.line_color)

    # 線の太さ
    if style.line_width:
        shape.line.width = Pt(style.line_width)

    # 線のスタイル
    if style.line_dash and style.line_dash in LINE_DASH_MAP:
        shape.line.dash_style = LINE_DASH_MAP[style.line_dash]


def _add_shape(
    slide: "Slide",
    content: ShapeContent,
) -> "Shape":
    """図形をスライドに追加"""
    shape_type = SHAPE_TYPE_MAP.get(content.shape_type, MSO_SHAPE.RECTANGLE)

    shape = slide.shapes.add_shape(
        shape_type,
        Inches(content.left),
        Inches(content.top),
        Inches(content.width),
        Inches(content.height),
    )

    # スタイルを適用
    _apply_shape_style(shape, content.style)

    # 回転
    if content.rotation:
        shape.rotation = content.rotation

    # テキストを追加
    if content.text and shape.has_text_frame:
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = content.text
        _apply_text_style(run, content.text_style)

    return shape


def _add_textbox(
    slide: "Slide",
    content: TextBoxContent,
) -> "Shape":
    """テキストボックスをスライドに追加"""
    textbox = slide.shapes.add_textbox(
        Inches(content.left),
        Inches(content.top),
        Inches(content.width),
        Inches(content.height),
    )

    text_frame = textbox.text_frame

    # 段落が指定されている場合
    if content.paragraphs:
        text_frame.clear()
        for i, para_content in enumerate(content.paragraphs):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = para_content.text
            _apply_text_style(run, para_content.style)
            if para_content.alignment and para_content.alignment in ALIGNMENT_MAP:
                p.alignment = ALIGNMENT_MAP[para_content.alignment]
    else:
        # 単純なテキスト
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = content.text
        _apply_text_style(run, content.style)

    # 背景色
    if content.fill_color:
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = _parse_rgb_color(content.fill_color)

    # 枠線色
    if content.line_color:
        textbox.line.color.rgb = _parse_rgb_color(content.line_color)

    return textbox


def _add_connector(
    slide: "Slide",
    content: ConnectorContent,
) -> "Shape":
    """コネクタ（接続線）をスライドに追加"""
    # 始点と終点から位置とサイズを計算
    start_x = Inches(content.start_x)
    start_y = Inches(content.start_y)
    end_x = Inches(content.end_x)
    end_y = Inches(content.end_y)

    # 直線コネクタを追加
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        start_x,
        start_y,
        end_x,
        end_y,
    )

    # 線の色
    if content.line_color:
        connector.line.color.rgb = _parse_rgb_color(content.line_color)

    # 線の太さ
    if content.line_width:
        connector.line.width = Pt(content.line_width)

    return connector


# =============================================================================
# メイン生成クラス
# =============================================================================

class PptxGenerator:
    """
    PowerPoint生成クラス
    
    Usage:
        generator = PptxGenerator(template_path, template_meta)
        generator.add_slide(slide_definition)
        generator.save(output_path)
    """
    
    def __init__(
        self,
        template_path: str | Path,
        template_meta: TemplateMeta | None = None,
    ):
        """
        Args:
            template_path: テンプレートPPTXのパス
            template_meta: テンプレートメタデータ（レイアウト検索に使用）
        """
        self.template_path = Path(template_path)
        self.template_meta = template_meta
        self.prs = Presentation(template_path)
        self.warnings: list[str] = []
        
        # テンプレートに既存スライドがあれば削除
        self._remove_existing_slides()
    
    def _remove_existing_slides(self) -> None:
        """テンプレートの既存スライドを削除"""
        # python-pptxには公式のスライド削除機能がないため、
        # 内部XMLを直接操作する
        xml_slides = self.prs.slides._sldIdLst
        slides_to_remove = list(xml_slides)
        for slide in slides_to_remove:
            xml_slides.remove(slide)
    
    def add_slide(self, definition: SlideDefinition) -> "Slide":
        """
        スライドを追加

        Args:
            definition: スライド定義

        Returns:
            追加されたスライド
        """
        # レイアウトを取得
        layout = _find_layout_by_name_or_index(
            self.prs,
            name=definition.layout_name,
            index=definition.layout_index,
            meta=self.template_meta,
        )

        # スライドを追加
        slide = self.prs.slides.add_slide(layout)

        # コンテンツを挿入
        for placeholder_key, content in definition.contents.items():
            self._insert_content(slide, placeholder_key, content)

        # 図形を追加
        for shape_content in definition.shapes:
            self._add_shape_to_slide(slide, shape_content)

        # スピーカーノート
        if definition.speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = definition.speaker_notes

        return slide

    def _add_shape_to_slide(
        self,
        slide: "Slide",
        content: ShapeContent | TextBoxContent | ConnectorContent,
    ) -> None:
        """スライドに図形を追加"""
        try:
            if isinstance(content, ShapeContent):
                _add_shape(slide, content)
            elif isinstance(content, TextBoxContent):
                _add_textbox(slide, content)
            elif isinstance(content, ConnectorContent):
                _add_connector(slide, content)
        except Exception as e:
            self.warnings.append(f"Failed to add shape: {str(e)}")
    
    def _insert_content(
        self,
        slide: "Slide",
        placeholder_key: str,
        content: SlideContent | str | list[str],
    ) -> None:
        """
        プレースホルダーにコンテンツを挿入
        
        Args:
            slide: スライド
            placeholder_key: プレースホルダーのキー（名前、インデックス、タイプ）
            content: 挿入するコンテンツ
        """
        # プレースホルダーを検索
        placeholder = self._resolve_placeholder(slide, placeholder_key)
        
        if placeholder is None:
            self.warnings.append(f"Placeholder not found: {placeholder_key}")
            return
        
        if not hasattr(placeholder, "text_frame"):
            self.warnings.append(f"Placeholder has no text frame: {placeholder_key}")
            return
        
        text_frame = placeholder.text_frame
        
        # コンテンツタイプに応じて挿入
        if isinstance(content, str):
            _insert_simple_text(text_frame, content)
        
        elif isinstance(content, list):
            # 文字列のリストは箇条書きとして扱う
            _insert_simple_bullets(text_frame, content)
        
        elif isinstance(content, TextContent):
            _insert_text_content(text_frame, content)
        
        elif isinstance(content, BulletContent):
            _insert_bullet_content(text_frame, content)
        
        elif isinstance(content, TableContent):
            _insert_table(slide, placeholder, content)
        
        elif isinstance(content, ChartContent):
            _insert_chart(slide, placeholder, content)
        
        elif isinstance(content, ImageContent):
            _insert_image(slide, placeholder, content)
    
    def _resolve_placeholder(
        self,
        slide: "Slide",
        key: str,
    ) -> "SlidePlaceholder | None":
        """
        キーからプレースホルダーを解決
        
        キーの形式:
        - "タイトル" → 名前で検索
        - "idx:0" → インデックスで検索
        - "type:TITLE" → タイプで検索
        """
        # インデックス指定
        if key.startswith("idx:"):
            idx = int(key[4:])
            return _find_placeholder(slide, idx=idx)
        
        # タイプ指定
        if key.startswith("type:"):
            type_str = key[5:]
            try:
                ph_type = PlaceholderType(type_str)
                return _find_placeholder(slide, ph_type=ph_type)
            except ValueError:
                pass
        
        # 名前で検索
        return _find_placeholder(slide, name=key)
    
    def save(self, output_path: str | Path) -> Path:
        """
        PPTXを保存
        
        Args:
            output_path: 出力パス
        
        Returns:
            保存されたファイルのパス
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        return output_path
    
    @property
    def slide_count(self) -> int:
        """スライド数を取得"""
        return len(self.prs.slides)


# =============================================================================
# 便利関数
# =============================================================================

def generate_pptx(
    template_path: str | Path,
    slides: list[SlideDefinition],
    output_path: str | Path,
    template_meta: TemplateMeta | None = None,
) -> tuple[Path, list[str]]:
    """
    PPTXを生成する便利関数
    
    Args:
        template_path: テンプレートPPTXのパス
        slides: スライド定義のリスト
        output_path: 出力パス
        template_meta: テンプレートメタデータ
    
    Returns:
        (出力パス, 警告リスト)
    """
    generator = PptxGenerator(template_path, template_meta)
    
    for slide_def in slides:
        generator.add_slide(slide_def)
    
    saved_path = generator.save(output_path)
    
    return saved_path, generator.warnings


def generate_from_ai_json(
    template_path: str | Path,
    ai_output: dict[str, Any],
    output_path: str | Path,
    template_meta: TemplateMeta | None = None,
) -> tuple[Path, list[str]]:
    """
    AIからのJSON出力を直接処理してPPTXを生成

    Args:
        template_path: テンプレートPPTXのパス
        ai_output: AIからのJSON出力（slides配列を含む）
        output_path: 出力パス
        template_meta: テンプレートメタデータ

    Returns:
        (出力パス, 警告リスト)
    """
    slides_data = ai_output.get("slides", [])

    slides = []
    for slide_data in slides_data:
        # AIからの出力をSlideDefinitionに変換
        slide_def = SlideDefinition(
            layout_name=slide_data.get("layoutName") or slide_data.get("layout_name"),
            layout_index=slide_data.get("layoutIndex") or slide_data.get("layout_index"),
            contents=_normalize_contents(slide_data.get("content", {})),
            shapes=_normalize_shapes(slide_data.get("shapes", [])),
            speaker_notes=slide_data.get("speakerNotes") or slide_data.get("speaker_notes"),
        )
        slides.append(slide_def)

    return generate_pptx(template_path, slides, output_path, template_meta)


def _normalize_shapes(
    shapes_data: list[dict[str, Any]],
) -> list[ShapeContent | TextBoxContent | ConnectorContent]:
    """
    AIからの図形データを正規化
    """
    normalized = []

    for shape_data in shapes_data:
        shape_type = shape_data.get("type", "shape")

        try:
            if shape_type == "shape":
                normalized.append(ShapeContent(**shape_data))
            elif shape_type == "textbox":
                normalized.append(TextBoxContent(**shape_data))
            elif shape_type == "connector":
                normalized.append(ConnectorContent(**shape_data))
        except Exception:
            # パース失敗時はスキップ
            continue

    return normalized


def _normalize_contents(
    contents: dict[str, Any],
) -> dict[str, SlideContent | str | list[str]]:
    """
    AIからのコンテンツをスキーマに正規化
    """
    normalized = {}
    
    for key, value in contents.items():
        if isinstance(value, str):
            normalized[key] = value
        elif isinstance(value, list):
            # 文字列のリストはそのまま
            if all(isinstance(v, str) for v in value):
                normalized[key] = value
            else:
                # 複雑な構造は TextContent に変換
                paragraphs = []
                for item in value:
                    if isinstance(item, str):
                        paragraphs.append(ParagraphContent(text=item))
                    elif isinstance(item, dict):
                        paragraphs.append(ParagraphContent(**item))
                normalized[key] = TextContent(paragraphs=paragraphs)
        elif isinstance(value, dict):
            # typeフィールドで判定
            content_type = value.get("type")
            if content_type == "text":
                normalized[key] = TextContent(**value)
            elif content_type == "bullets":
                normalized[key] = BulletContent(**value)
            elif content_type == "table":
                normalized[key] = TableContent(**value)
            elif content_type == "chart":
                normalized[key] = ChartContent(**value)
            elif content_type == "image":
                normalized[key] = ImageContent(**value)
            else:
                # typeがない場合は paragraphs の有無で判定
                if "paragraphs" in value:
                    normalized[key] = TextContent(**value)
                else:
                    # 単純なkey-valueとして処理
                    normalized[key] = str(value)
    
    return normalized
