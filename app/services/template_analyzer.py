"""
テンプレート解析サービス
PPTXファイルからレイアウト・プレースホルダー情報を抽出
"""
from __future__ import annotations

import hashlib
import json
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

from app.schemas import (
    LayoutMeta,
    PlaceholderMeta,
    PlaceholderType,
    TemplateCustomConfig,
    TemplateMeta,
)

if TYPE_CHECKING:
    from pptx.slide import SlideLayout
    from pptx.shapes.placeholder import LayoutPlaceholder


# プレースホルダータイプのマッピング
PLACEHOLDER_TYPE_MAP: dict[int | None, PlaceholderType | None] = {
    PP_PLACEHOLDER.TITLE: PlaceholderType.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE: PlaceholderType.CENTER_TITLE,
    PP_PLACEHOLDER.SUBTITLE: PlaceholderType.SUBTITLE,
    PP_PLACEHOLDER.BODY: PlaceholderType.BODY,
    PP_PLACEHOLDER.OBJECT: PlaceholderType.OBJECT,
    PP_PLACEHOLDER.CHART: PlaceholderType.CHART,
    PP_PLACEHOLDER.TABLE: PlaceholderType.TABLE,
    PP_PLACEHOLDER.PICTURE: PlaceholderType.PICTURE,
    PP_PLACEHOLDER.FOOTER: PlaceholderType.FOOTER,
    PP_PLACEHOLDER.DATE: PlaceholderType.DATE,
    PP_PLACEHOLDER.SLIDE_NUMBER: PlaceholderType.SLIDE_NUMBER,
}


def _emu_to_inches(emu: int) -> float:
    """EMUをインチに変換"""
    return round(emu / 914400, 2)


def _get_placeholder_type(ph_type: int | None) -> PlaceholderType | None:
    """python-pptxのプレースホルダータイプを変換"""
    if ph_type is None:
        return None
    return PLACEHOLDER_TYPE_MAP.get(ph_type)


def _extract_placeholder_meta(placeholder: "LayoutPlaceholder", idx: int) -> PlaceholderMeta:
    """プレースホルダーからメタデータを抽出"""
    ph_type = None
    if hasattr(placeholder, "placeholder_format") and placeholder.placeholder_format:
        ph_type = _get_placeholder_type(placeholder.placeholder_format.type)
    
    # プレースホルダー名を取得
    name = placeholder.name or f"Placeholder {idx}"
    
    # テキストフレームの有無を確認
    has_text_frame = hasattr(placeholder, "has_text_frame") and placeholder.has_text_frame
    
    return PlaceholderMeta(
        idx=placeholder.placeholder_format.idx if placeholder.placeholder_format else idx,
        type=ph_type,
        name=name,
        left=_emu_to_inches(placeholder.left),
        top=_emu_to_inches(placeholder.top),
        width=_emu_to_inches(placeholder.width),
        height=_emu_to_inches(placeholder.height),
        has_text_frame=has_text_frame,
    )


def _extract_layout_meta(layout: "SlideLayout", index: int) -> LayoutMeta:
    """レイアウトからメタデータを抽出"""
    placeholders: list[PlaceholderMeta] = []
    
    for idx, shape in enumerate(layout.placeholders):
        # スライド番号、日付、フッターは除外（通常コンテンツ用ではない）
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            ph_type = shape.placeholder_format.type
            if ph_type in (
                PP_PLACEHOLDER.SLIDE_NUMBER,
                PP_PLACEHOLDER.DATE,
                PP_PLACEHOLDER.FOOTER,
            ):
                continue
        
        placeholder_meta = _extract_placeholder_meta(shape, idx)
        placeholders.append(placeholder_meta)
    
    # 位置でソート（上から下、左から右）
    placeholders.sort(key=lambda p: (p.top, p.left))
    
    return LayoutMeta(
        index=index,
        name=layout.name or f"Layout {index}",
        placeholders=placeholders,
    )


def analyze_template(
    pptx_path: str | Path,
    template_id: str | None = None,
    custom_config: TemplateCustomConfig | None = None,
) -> TemplateMeta:
    """
    テンプレートPPTXを解析してメタデータを生成

    Args:
        pptx_path: PPTXファイルのパス
        template_id: テンプレートID（指定しない場合はハッシュから生成）
        custom_config: カスタム設定

    Returns:
        TemplateMeta: テンプレートのメタデータ
    """
    pptx_path = Path(pptx_path)

    if not pptx_path.exists():
        raise FileNotFoundError(f"Template file not found: {pptx_path}")

    # テンプレートIDを生成（指定がなければファイルハッシュから）
    if template_id is None:
        file_hash = hashlib.md5(pptx_path.read_bytes()).hexdigest()[:8]
        template_id = f"{pptx_path.stem}_{file_hash}"

    prs = Presentation(pptx_path)

    # デバッグ: スライドマスター数を出力
    print(f"[DEBUG] Template: {pptx_path}")
    print(f"[DEBUG] Number of slide masters: {len(prs.slide_masters)}")

    # レイアウト情報を抽出（全スライドマスターから）
    # prs.slide_layouts は最初のスライドマスターのみを返すため、
    # 全てのスライドマスターをイテレートして全レイアウトを取得する
    layouts: list[LayoutMeta] = []
    global_index = 0

    for master_idx, master in enumerate(prs.slide_masters):
        master_layouts = list(master.slide_layouts)
        print(f"[DEBUG] Master {master_idx}: {len(master_layouts)} layouts")
        for layout in master_layouts:
            print(f"[DEBUG]   Layout [{global_index}]: {layout.name}")
            layout_meta = _extract_layout_meta(layout, global_index)
            layouts.append(layout_meta)
            global_index += 1

    print(f"[DEBUG] Total layouts found: {len(layouts)}")

    return TemplateMeta(
        id=template_id,
        name=pptx_path.stem,
        file_name=pptx_path.name,
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        layouts=layouts,
        custom_config=custom_config,
    )


def enrich_template_meta(
    meta: TemplateMeta,
    layout_descriptions: dict[str, str] | None = None,
    layout_hints: dict[str, str] | None = None,
    layout_aliases: dict[str, int] | None = None,
) -> TemplateMeta:
    """
    テンプレートメタデータにAI向け説明を追加
    
    Args:
        meta: 元のメタデータ
        layout_descriptions: レイアウト名 → 説明のマッピング
        layout_hints: レイアウト名 → 使用ヒントのマッピング
        layout_aliases: エイリアス → レイアウトインデックスのマッピング
    
    Returns:
        TemplateMeta: 説明が追加されたメタデータ
    """
    layout_descriptions = layout_descriptions or {}
    layout_hints = layout_hints or {}
    
    # レイアウトに説明を追加
    enriched_layouts = []
    for layout in meta.layouts:
        enriched_layout = layout.model_copy(deep=True)
        
        if layout.name in layout_descriptions:
            enriched_layout.description = layout_descriptions[layout.name]
        
        if layout.name in layout_hints:
            enriched_layout.usage_hint = layout_hints[layout.name]
        
        enriched_layouts.append(enriched_layout)
    
    # カスタム設定を更新
    custom_config = meta.custom_config or TemplateCustomConfig()
    if layout_aliases:
        custom_config.layout_aliases = layout_aliases
    
    return meta.model_copy(
        update={
            "layouts": enriched_layouts,
            "custom_config": custom_config,
        }
    )


def save_template_meta(meta: TemplateMeta, output_path: str | Path) -> None:
    """メタデータをJSONファイルに保存"""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(meta.model_dump(), f, ensure_ascii=False, indent=2)


def load_template_meta(meta_path: str | Path) -> TemplateMeta:
    """JSONファイルからメタデータを読み込み"""
    meta_path = Path(meta_path)
    
    if not meta_path.exists():
        raise FileNotFoundError(f"Meta file not found: {meta_path}")
    
    with open(meta_path, encoding="utf-8") as f:
        data = json.load(f)
    
    return TemplateMeta.model_validate(data)


# =============================================================================
# デフォルトのレイアウト説明（日本語PowerPointの標準レイアウト用）
# =============================================================================

DEFAULT_LAYOUT_DESCRIPTIONS: dict[str, str] = {
    "タイトル スライド": "プレゼンテーションの表紙。タイトルとサブタイトル（日付・発表者など）を配置",
    "タイトルとコンテンツ": "最も汎用的なレイアウト。タイトルと本文（箇条書き可）",
    "セクション見出し": "章・セクションの区切りに使用。大きなタイトルと補足テキスト",
    "2 つのコンテンツ": "左右2カラムで比較や並列表示に最適",
    "比較": "2つの項目を明確に比較するレイアウト。各カラムに見出し付き",
    "タイトルのみ": "タイトルのみ。本文エリアは自由に使用可能",
    "白紙": "プレースホルダーなし。完全に自由なレイアウト",
    "タイトル付きのコンテンツ": "メインコンテンツとタイトルの標準構成",
    "タイトル付きの図": "画像や図を大きく配置するレイアウト",
}

DEFAULT_LAYOUT_HINTS: dict[str, str] = {
    "タイトル スライド": "プレゼンの最初のスライドに使用",
    "タイトルとコンテンツ": "説明、リスト、概要など幅広く使用可能",
    "セクション見出し": "トピックが変わる際の区切りに使用",
    "2 つのコンテンツ": "Before/After、メリット/デメリット、2案の比較に最適",
    "比較": "機能比較、製品比較などに使用",
    "白紙": "カスタムデザインや図表を自由配置する場合に使用",
}
